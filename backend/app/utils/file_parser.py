"""
Best-effort text extraction from common admin-upload file types.

The Bedrock KB connector parses the file itself during ingestion — this
extractor is only used for:
  * showing a preview / first-N-chars in the admin UI
  * an offline OCR / search backup if the KB is unavailable
"""
from __future__ import annotations

import io
from pathlib import Path

ALLOWED_EXT = {".pdf", ".docx", ".xlsx", ".csv", ".pptx", ".ppt", ".txt"}


def is_supported(filename: str) -> bool:
    return Path(filename).suffix.lower() in ALLOWED_EXT


def extract_text(filename: str, data: bytes, *, max_chars: int = 4000) -> str:
    ext = Path(filename).suffix.lower()
    try:
        if ext == ".txt" or ext == ".csv":
            return data.decode("utf-8", errors="ignore")[:max_chars]
        if ext == ".pdf":
            from pypdf import PdfReader
            r = PdfReader(io.BytesIO(data))
            chunks = []
            for page in r.pages[:20]:
                chunks.append(page.extract_text() or "")
                if sum(len(c) for c in chunks) > max_chars:
                    break
            return "\n".join(chunks)[:max_chars]
        if ext == ".docx":
            from docx import Document
            doc = Document(io.BytesIO(data))
            return "\n".join(p.text for p in doc.paragraphs)[:max_chars]
        if ext == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            out = []
            for sheet in wb.sheetnames[:5]:
                out.append(f"[Sheet: {sheet}]")
                ws = wb[sheet]
                for i, row in enumerate(ws.iter_rows(values_only=True)):
                    if i > 100:
                        break
                    out.append(",".join("" if c is None else str(c) for c in row))
            return "\n".join(out)[:max_chars]
        if ext in (".pptx", ".ppt"):
            from pptx import Presentation
            pres = Presentation(io.BytesIO(data))
            chunks = []
            for slide in pres.slides[:30]:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        chunks.append(shape.text)
            return "\n".join(chunks)[:max_chars]
    except Exception:
        return ""
    return ""
